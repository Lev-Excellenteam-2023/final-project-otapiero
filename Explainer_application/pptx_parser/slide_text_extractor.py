"""
 this file is used to extract text from pptx file
"""

import pptx

from Explainer_application.pptx_parser import text_processing, pptx_reader


class SlideTextExtractor:
    """
    This class is used to extract text from a pptx file.


    """
    def __init__(self, pptx_file_path: str ):
        """
        initialize the class.

        :param pptx_file_path: str, path to the pptx file.
        """
        self.pptx_file_path = pptx_file_path
        self.pptx_file_slides = pptx_reader.read_pptx_file(pptx_file_path)
        self.index = 0



    def extract_slide( self, slide: pptx.slide.Slide ) ->str:
        """
        extract text from a slide in the pptx file.

        :return: list of lists, each list contains title, text and notes of a slide.
        """
        slide_title = self.get_slide_title(slide)
        slide_text = self.get_slide_text(slide)
        slide_notes = self.get_slide_notes(slide)
        slide_content = text_processing.process_pptx_slide_content([slide_title, slide_text, slide_notes])
        return slide_content


    def get_slide_text(self, slide: pptx.slide.Slide) -> str:
        """
        extract text from a slide.

        :param slide: pptx.slide.Slide object.
        :return: str, text from the slide.
        """
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
        return text

    def get_slide_title(self, slide: pptx.slide.Slide ) -> str:
        """
        extract title from a slide.

        :param slide: pptx.slide.Slide object.
        :return: str, title from the slide
        """
        title = slide.shapes.title
        if title is not None:
            return title.text
        return ""

    def get_slide_notes(self, slide: pptx.slide.Slide ) -> str:
        """
        extract notes from a slide.

        :param slide: pptx.slide.Slide object.
        :return: str, notes from the slide.
        """
        notes = ""
        for shape in slide.notes_slide.shapes:
            if hasattr(shape, "text"):
                notes += shape.text
        return notes

    def __iter__(self):
        """
        iterate over the slides in the pptx file.
        """
        self.index = 0
        return self

    def __next__(self):
        """
        get the next slide in the pptx file.
        """
        if self.index >= len(self.pptx_file_slides):
            raise StopIteration
        else:
            slide = self.pptx_file_slides[self.index]
            self.index += 1
            return self.extract_slide(slide)











if __name__ == "__main__":
    pptx_file_path = str(input("Enter PowerPoint file path:"))
    if pptx_file_path == "":
        pptx_file_path = r"C:\Users\ouriel\Desktop\End of course exercise - kickof - upload.pptx"
    slide_text_extractor = SlideTextExtractor(pptx_file_path=pptx_file_path)
    for slide in slide_text_extractor:
        print(slide)
        print("---------------------------------------------------------------------")