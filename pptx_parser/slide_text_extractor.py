"""
 this file is used to extract text from pptx file
"""
import pptx

from pptx_parser import pptx_reader
from pptx_parser import text_processing


def get_pptx_file_content( pptx_file_path: str ) -> list[list[str]]:
    """
    extract text from a pptx file.

    :param pptx_file_path: str, path to the pptx file.
    :return list of lists, each list contains title, text and notes of a slide.
    """
    pptx_file_slides = pptx_reader.read_pptx_file(pptx_file_path)
    pptx_file_content = []
    for slide in pptx_file_slides:
        slide_title = get_slide_title(slide)
        slide_text = get_slide_text(slide)
        slide_notes = get_slide_notes(slide)
        pptx_file_content.append([slide_title, slide_text, slide_notes])
    pptx_file_content = text_processing.process_pptx_file_content(pptx_file_content)
    return pptx_file_content



def get_slide_text ( slide: pptx.slide.Slide ) -> str:
    """
    extract text from a slide.

    :param slide: pptx.slide.Slide object.
    :return: str, text from the slide.
    """
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text
    return text


def get_slide_title ( slide: pptx.slide.Slide ) -> str:
    """
    extract title from a slide.

    :param slide: pptx.slide.Slide object.
    :return: str, title from the slide
    """
    title = slide.shapes.title
    if title is not None:
        return title.text
    return ""


def get_slide_notes ( slide: pptx.slide.Slide ) -> str:
    """
    extract notes from a slide.

    :param slide: pptx.slide.Slide object.
    :return: str, notes from the slide.
    """
    notes = ""
    for shape in slide.notes_slide.shapes:
        if hasattr(shape, "text"):
            notes += shape.text
    return notes


if __name__ == "__main__":
    pptx_file_path = str(input("Enter PowerPoint file path:"))
    if pptx_file_path == "":
        pptx_file_path = r"C:\Users\ouriel\Desktop\End of course exercise - kickof - upload.pptx"
    pptx_file_content = get_pptx_file_content(pptx_file_path)
    for slide in pptx_file_content:
        print(slide)
        print("---------------------------------------------------------------------")
